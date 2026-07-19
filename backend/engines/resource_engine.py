import json
import io
import uuid
import base64
from typing import Dict, List, Any, Optional
from utils.llm_client import invoke_llm
from config import settings
import logging

logger = logging.getLogger(__name__)


class ResourceGenerator:
    def __init__(self):
        pass

    def generate_ppt(self, topic: str, content: str) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = topic
            subtitle.text = "青藜伴读"

            prompt = f"""你是一个专业的PPT设计师。请为主题"{topic}"生成PPT内容。
要求：
1. 生成6-8个幻灯片
2. 每个幻灯片包含标题和3-5个要点
3. 内容基于以下知识：{content}
4. 返回纯JSON格式，不要任何解释"""

            response = invoke_llm(prompt)
            try:
                ppt_data = json.loads(response)
            except:
                ppt_data = self._generate_default_ppt_content(topic, content)

            for slide_info in ppt_data.get("slides", []):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                title_shape = slide.shapes.title
                title_shape.text = slide_info.get("title", "")
                
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                
                for point in slide_info.get("points", []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(50, 50, 50)

            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)

            return {
                "type": "ppt",
                "status": "success",
                "data": base64.b64encode(ppt_buffer.read()).decode(),
                "filename": f"{topic}.pptx",
                "slide_count": len(prs.slides),
                "topic": topic,
            }
        except ImportError:
            return self._generate_ppt_fallback(topic, content)
        except Exception as e:
            logger.error(f"PPT生成失败: {e}")
            return self._generate_ppt_fallback(topic, content)

    def _generate_ppt_fallback(self, topic: str, content: str) -> Dict[str, Any]:
        prompt = f"""请为主题"{topic}"生成PPT大纲，包含6-8个幻灯片的标题和要点。
基于内容：{content}
返回纯JSON格式。"""
        response = invoke_llm(prompt)
        try:
            ppt_data = json.loads(response)
        except:
            ppt_data = self._generate_default_ppt_content(topic, content)
        
        return {
            "type": "ppt",
            "status": "success",
            "data": json.dumps(ppt_data, ensure_ascii=False),
            "format": "json",
            "message": "python-pptx未安装，返回PPT大纲JSON",
            "topic": topic,
        }

    def _generate_default_ppt_content(self, topic: str, content: str) -> Dict[str, Any]:
        lines = content.split('\n')[:20]
        return {
            "slides": [
                {"title": "目录", "points": ["概述", "核心概念", "关键原理", "实践应用", "总结"]},
                {"title": "概述", "points": [topic, "学习目标", "重要性"]},
                {"title": "核心概念", "points": lines[:5]},
                {"title": "关键原理", "points": lines[5:10]},
                {"title": "实践应用", "points": lines[10:15]},
                {"title": "总结", "points": ["要点回顾", "学习建议", "后续计划"]},
            ]
        }

    def generate_mindmap(self, topic: str, content: str) -> Dict[str, Any]:
        prompt = f"""你是一个思维导图专家。请为主题"{topic}"生成思维导图结构。
要求：
1. 返回JSON格式，包含nodes和edges
2. nodes: [{{"id": "...", "label": "...", "level": 0/1/2/3}}]
3. edges: [{{"from": "...", "to": "..."}}]
4. 包含3-4层节点
5. 内容基于以下知识：{content}
6. 返回纯JSON，不要任何解释"""

        response = invoke_llm(prompt)
        try:
            mindmap_data = json.loads(response)
        except:
            mindmap_data = self._generate_default_mindmap(topic, content)

        svg_content = self._render_mindmap_to_svg(mindmap_data, topic)

        return {
            "type": "mindmap",
            "status": "success",
            "svg": svg_content,
            "topic": topic,
        }

    def _generate_default_mindmap(self, topic: str, content: str) -> Dict[str, Any]:
        lines = content.split('\n')[:15]
        nodes = [
            {"id": "root", "label": topic, "level": 0},
            {"id": "c1", "label": "核心概念", "level": 1},
            {"id": "c2", "label": "关键原理", "level": 1},
            {"id": "c3", "label": "实践应用", "level": 1},
        ]
        edges = [
            {"from": "root", "to": "c1"},
            {"from": "root", "to": "c2"},
            {"from": "root", "to": "c3"},
        ]
        for i, line in enumerate(lines[:6]):
            node_id = f"n{i}"
            nodes.append({"id": node_id, "label": line[:20], "level": 2})
            edges.append({"from": ["c1", "c2", "c3"][i % 3], "to": node_id})
        return {"nodes": nodes, "edges": edges}

    def _render_mindmap_to_svg(self, mindmap_data: Dict[str, Any], topic: str) -> str:
        try:
            import pydot
            graph = pydot.Dot(graph_type='digraph', rankdir='LR', bgcolor='transparent')
            
            nodes = mindmap_data.get("nodes", [])
            edges = mindmap_data.get("edges", [])
            
            color_map = {0: "#4A90D9", 1: "#67C23A", 2: "#E6A23C", 3: "#F56C6C"}
            
            for node in nodes:
                node_id = node.get("id", "")
                label = node.get("label", "")
                level = node.get("level", 0)
                color = color_map.get(level, "#909399")
                
                graph.add_node(pydot.Node(
                    node_id,
                    label=label,
                    shape="box",
                    style="filled",
                    fillcolor=color,
                    fontcolor="white",
                    fontsize="14",
                    fontname="Microsoft YaHei",
                    padding="0.5",
                ))
            
            for edge in edges:
                graph.add_edge(pydot.Edge(edge.get("from"), edge.get("to")))
            
            svg_bytes = graph.create_svg()
            return svg_bytes.decode('utf-8')
        except Exception as e:
            logger.warning(f"SVG生成失败，使用HTML替代: {e}")
            return self._generate_html_mindmap(mindmap_data, topic)

    def _generate_html_mindmap(self, mindmap_data: Dict[str, Any], topic: str) -> str:
        nodes = mindmap_data.get("nodes", [])
        edges = mindmap_data.get("edges", [])
        
        html = f'<div class="mindmap-container"><h3>{topic}</h3><ul class="mindmap">'
        root_nodes = [n for n in nodes if n.get("level") == 0]
        
        for root in root_nodes:
            html += f'<li><strong>{root.get("label")}</strong>'
            level1_nodes = [n for n in nodes if any(e.get("from") == root.get("id") and e.get("to") == n.get("id") for e in edges)]
            
            if level1_nodes:
                html += '<ul>'
                for n1 in level1_nodes:
                    html += f'<li>{n1.get("label")}'
                    level2_nodes = [n for n in nodes if any(e.get("from") == n1.get("id") and e.get("to") == n.get("id") for e in edges)]
                    
                    if level2_nodes:
                        html += '<ul>'
                        for n2 in level2_nodes:
                            html += f'<li>{n2.get("label")}</li>'
                        html += '</ul>'
                    html += '</li>'
                html += '</ul>'
            html += '</li>'
        
        html += '</ul></div>'
        return html

    def generate_code(self, topic: str, content: str) -> Dict[str, Any]:
        prompt = f"""你是一个编程专家。请为主题"{topic}"生成可运行的Python代码示例。
要求：
1. 生成完整的、可直接运行的代码
2. 包含注释说明
3. 包含测试用例
4. 代码要能体现"{topic}"的核心概念
5. 基于以下知识：{content}
6. 返回格式：{{"code": "...", "tests": "...", "explanation": "..."}}"""

        response = invoke_llm(prompt)
        try:
            code_data = json.loads(response)
        except:
            code_data = {
                "code": f"# {topic} - 代码示例\n# 基于知识内容：{content[:100]}...\n\ndef main():\n    print(f'{topic} 代码示例')\n\nif __name__ == '__main__':\n    main()",
                "tests": "# 测试用例\n# 运行方式: python script.py\n# 预期输出: 正常运行无报错",
                "explanation": f"{topic}的代码实现示例",
            }

        return {
            "type": "code",
            "status": "success",
            "code": code_data.get("code", ""),
            "tests": code_data.get("tests", ""),
            "explanation": code_data.get("explanation", ""),
            "topic": topic,
            "language": "python",
        }

    def generate_reading(self, topic: str, content: str) -> Dict[str, Any]:
        try:
            from rag.engine import RAGEngine
            rag_engine = RAGEngine()
            results = rag_engine.retrieve(topic, top_k=5)
            
            contexts = [r.get("text", "") for r in results]
            context_text = "\n\n".join(contexts)

            prompt = f"""你是一个专业的文献综述专家。请基于以下检索到的资料，为主题"{topic}"生成拓展阅读内容。
检索资料：
{context_text}

要求：
1. 生成一篇500-800字的综述文章
2. 包含核心概念、发展历程、研究现状、未来趋势
3. 列出关键参考文献
4. 返回格式：{{"summary": "...", "references": [...]}}"""

            response = invoke_llm(prompt)
            try:
                reading_data = json.loads(response)
            except:
                reading_data = {
                    "summary": f"## {topic} 拓展阅读\n\n{content[:500]}...",
                    "references": ["检索到的相关文档"],
                }

            reading_data["sources"] = results[:3]

            return {
                "type": "reading",
                "status": "success",
                "summary": reading_data.get("summary", ""),
                "references": reading_data.get("references", []),
                "sources": reading_data.get("sources", []),
                "topic": topic,
            }
        except Exception as e:
            logger.warning(f"RAG检索失败，使用内容摘要: {e}")
            
            prompt = f"""请为主题"{topic}"生成拓展阅读内容，基于以下知识：{content}
要求：
1. 生成一篇500-800字的综述文章
2. 包含核心概念、发展历程、研究现状、未来趋势
3. 列出推荐阅读材料"""

            response = invoke_llm(prompt)
            return {
                "type": "reading",
                "status": "success",
                "summary": response,
                "references": ["知识库内容"],
                "topic": topic,
                "message": "RAG引擎不可用，使用知识库内容生成",
            }

    def generate_video(self, topic: str, content: str) -> Dict[str, Any]:
        prompt = f"""你是一个视频脚本专家。请为主题"{topic}"生成一个短视频脚本。
要求：
1. 视频时长3-5分钟
2. 包含分镜脚本、旁白文案、画面描述
3. 基于以下知识：{content}
4. 返回格式：{{"script": [...], "voiceover": "...", "duration": 180}}"""

        response = invoke_llm(prompt)
        try:
            video_data = json.loads(response)
        except:
            video_data = {
                "script": [
                    {"time": "0-30s", "scene": f"{topic}介绍", "description": "主题画面"},
                    {"time": "30-90s", "scene": "核心概念", "description": "图文展示"},
                    {"time": "90-150s", "scene": "实际应用", "description": "动画演示"},
                    {"time": "150-180s", "scene": "总结", "description": "要点回顾"},
                ],
                "voiceover": f"{topic}是重要的知识点...",
                "duration": 180,
            }

        try:
            from utils.xfyun_tts import xfyun_synthesize_speech
            from utils.tts_client import synthesize_speech
            
            voiceover_text = video_data.get("voiceover", "")
            audio_data = None
            
            if settings.XFYUN_TTS_APP_ID:
                audio_data = xfyun_synthesize_speech(voiceover_text, "default")
            
            if audio_data is None:
                audio_data = synthesize_speech(voiceover_text, "default")
            
            if audio_data:
                video_data["audio_base64"] = base64.b64encode(audio_data).decode()
                video_data["audio_available"] = True
        except Exception as e:
            logger.warning(f"TTS生成失败: {e}")
            video_data["audio_available"] = False

        return {
            "type": "video",
            "status": "success",
            "script": video_data.get("script", []),
            "voiceover": video_data.get("voiceover", ""),
            "duration": video_data.get("duration", 180),
            "audio_available": video_data.get("audio_available", False),
            "audio_base64": video_data.get("audio_base64", ""),
            "topic": topic,
        }


resource_generator = ResourceGenerator()
